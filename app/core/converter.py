"""
转换引擎 - 图片/音频/视频/PDF格式转换
"""
import os
import subprocess
from pathlib import Path
from typing import Optional, Callable

from PIL import Image

try:
    import fitz  # PyMuPDF
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import ffmpeg
    HAS_FFMPEG = True
except ImportError:
    HAS_FFMPEG = False

try:
    from pdf2docx import Converter as PdfDocxConverter
    HAS_PDF2DOCX = True
except ImportError:
    HAS_PDF2DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openpyxl import Workbook, load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


# ── 图片转换 ─────────────────────────────────────────────────

IMAGE_FORMATS = {
    'jpg': 'JPEG',
    'jpeg': 'JPEG',
    'png': 'PNG',
    'webp': 'WEBP',
    'bmp': 'BMP',
    'gif': 'GIF',
    'tiff': 'TIFF',
    'ico': 'ICO',
    'svg': 'SVG',  # 只支持SVG→其他，不支持其他→SVG
}


def get_image_formats() -> list[str]:
    return sorted(k for k in IMAGE_FORMATS if k != 'svg') + ['svg']


def convert_image(
    input_path: str,
    output_path: str,
    quality: int = 90,
    resize: Optional[tuple[int, int]] = None,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """图片格式转换"""
    try:
        img = Image.open(input_path)

        # 调整尺寸
        if resize:
            img = img.resize(resize, Image.LANCZOS)

        # 转换并保存
        output_ext = Path(output_path).suffix.lower().lstrip('.')
        save_kwargs = {}

        if output_ext in ('jpg', 'jpeg'):
            save_kwargs['quality'] = quality
            if img.mode in ('RGBA', 'P'):
                # JPEG不支持透明，转RGB
                background = Image.new('RGB', img.size, (255, 255, 255))
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                img = background
        elif output_ext == 'png':
            save_kwargs['compress_level'] = 6
        elif output_ext == 'webp':
            save_kwargs['quality'] = quality
        elif output_ext == 'ico':
            # ICO需要特定尺寸，不支持quality等参数
            if resize is None:
                img = img.resize((64, 64), Image.LANCZOS)
            save_kwargs = {}  # ICO 不支持额外参数

        img.save(output_path, **save_kwargs)

        if progress_callback:
            progress_callback(1, 1)

        return True

    except Exception as e:
        raise RuntimeError(f"图片转换失败: {e}")


def convert_images_batch(
    input_paths: list[str],
    output_dir: str,
    target_format: str,
    quality: int = 90,
    resize: Optional[tuple[int, int]] = None,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> int:
    """批量图片转换"""
    os.makedirs(output_dir, exist_ok=True)
    total = len(input_paths)
    success = 0
    target_format = target_format.lstrip('.')

    for i, path in enumerate(input_paths):
        try:
            stem = Path(path).stem
            out_path = os.path.join(output_dir, f"{stem}.{target_format}")
            convert_image(path, out_path, quality, resize)
            success += 1
        except Exception:
            continue

        if progress_callback:
            progress_callback(i + 1, total)

    return success


# ── 音频转换 ─────────────────────────────────────────────────

AUDIO_FORMATS = ['mp3', 'wav', 'flac', 'aac', 'ogg', 'm4a', 'wma']


def get_audio_formats() -> list[str]:
    return AUDIO_FORMATS


def convert_audio(
    input_path: str,
    output_path: str,
    bitrate: Optional[str] = None,  # e.g. '192k'
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """音频格式转换 (依赖 ffmpeg)"""
    if not HAS_FFMPEG:
        raise RuntimeError("音频转换需要 ffmpeg，请确保已安装")

    try:
        # 确保输出目录存在
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

        stream = ffmpeg.input(input_path)
        kwargs = {}
        if bitrate:
            kwargs['b:a'] = bitrate

        stream = ffmpeg.output(stream, output_path, **kwargs)
        ffmpeg.run(stream, overwrite_output=True, capture_stdout=True, capture_stderr=True)

        if progress_callback:
            progress_callback(1, 1)

        return True

    except Exception as e:
        raise RuntimeError(f"音频转换失败: {e}")


# ── 视频转换 ─────────────────────────────────────────────────

VIDEO_FORMATS = ['mp4', 'avi', 'mkv', 'mov', 'wmv', 'webm', 'flv']


def get_video_formats() -> list[str]:
    return VIDEO_FORMATS


def convert_video(
    input_path: str,
    output_path: str,
    video_bitrate: Optional[str] = None,
    audio_bitrate: Optional[str] = None,
    resolution: Optional[tuple[int, int]] = None,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """视频格式转换 (依赖 ffmpeg)"""
    if not HAS_FFMPEG:
        raise RuntimeError("视频转换需要 ffmpeg，请确保已安装")

    try:
        # 确保输出目录存在
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

        stream = ffmpeg.input(input_path)
        kwargs = {}

        if video_bitrate:
            kwargs['b:v'] = video_bitrate
        if audio_bitrate:
            kwargs['b:a'] = audio_bitrate
        if resolution:
            kwargs['vf'] = f'scale={resolution[0]}:{resolution[1]}'

        stream = ffmpeg.output(stream, output_path, **kwargs)
        ffmpeg.run(stream, overwrite_output=True, capture_stdout=True, capture_stderr=True)

        if progress_callback:
            progress_callback(1, 1)

        return True

    except Exception as e:
        raise RuntimeError(f"视频转换失败: {e}")


# ── PDF 转换 ─────────────────────────────────────────────────

def images_to_pdf(
    image_paths: list[str],
    output_path: str,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """多张图片合成 PDF（使用 Pillow 逐页方式）"""
    if not image_paths:
        raise RuntimeError("没有图片可供合成")

    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

    try:
        total = len(image_paths)
        rgb_images = []

        for i, path in enumerate(image_paths):
            img = Image.open(path)
            if img.mode != 'RGB':
                if img.mode == 'RGBA':
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[3])
                    img = background
                else:
                    img = img.convert('RGB')
            rgb_images.append(img)

            if progress_callback:
                progress_callback(i + 1, total)

        # Pillow 方式：直接用 save_all + append_images
        rgb_images[0].save(output_path, save_all=True, append_images=rgb_images[1:])
        return True

    except Exception as e:
        raise RuntimeError(f"图片转PDF失败: {e}")


def pdf_to_images(
    pdf_path: str,
    output_dir: str,
    image_format: str = 'png',
    dpi: int = 200,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> int:
    """PDF 每页转成图片"""
    if not HAS_PDF:
        raise RuntimeError("PDF 功能需要 PyMuPDF 库")

    os.makedirs(output_dir, exist_ok=True)

    try:
        doc = fitz.open(pdf_path)
        total = len(doc)
        success = 0

        for i in range(total):
            page = doc[i]
            pix = page.get_pixmap(dpi=dpi)
            out_path = os.path.join(output_dir, f"page_{i+1:03d}.{image_format}")
            pix.save(out_path)
            success += 1

            if progress_callback:
                progress_callback(i + 1, total)

        doc.close()
        return success

    except Exception as e:
        raise RuntimeError(f"PDF转图片失败: {e}")


# ── PDF ↔ Office 转换 ─────────────────────────────────────

def pdf_to_docx(
    pdf_path: str,
    output_path: str,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """PDF 转 Word (docx)"""
    if not HAS_PDF2DOCX:
        raise RuntimeError("PDF转Word需要 pdf2docx 库，请执行: pip install pdf2docx")

    try:
        cv = PdfDocxConverter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

        if progress_callback:
            progress_callback(1, 1)

        return True

    except Exception as e:
        raise RuntimeError(f"PDF转Word失败: {e}")


def pdf_to_pptx(
    pdf_path: str,
    output_path: str,
    image_dpi: int = 150,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """PDF 转 PPT (pptx) - 每页截图放入幻灯片"""
    if not HAS_PDF:
        raise RuntimeError("PDF 功能需要 PyMuPDF 库")
    if not HAS_PPTX:
        raise RuntimeError("PDF转PPT需要 python-pptx 库，请执行: pip install python-pptx")

    try:
        doc = fitz.open(pdf_path)
        total = len(doc)
        prs = Presentation()
        prs.slide_width = 12192000   # 16:9 宽屏
        prs.slide_height = 6858000

        import io

        for i in range(total):
            page = doc[i]
            pix = page.get_pixmap(dpi=image_dpi)
            img_bytes = pix.tobytes("png")

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            image_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(image_stream, 0, 0,
                                     width=prs.slide_width,
                                     height=prs.slide_height)

            if progress_callback:
                progress_callback(i + 1, total)

        doc.close()
        prs.save(output_path)
        return True

    except Exception as e:
        raise RuntimeError(f"PDF转PPT失败: {e}")


def pdf_to_excel(
    pdf_path: str,
    output_path: str,
    progress_callback: Optional[Callable[[int, int], None]] = None,
) -> bool:
    """PDF 转 Excel (xlsx) - 提取文本到单元格"""
    if not HAS_PDF:
        raise RuntimeError("PDF 功能需要 PyMuPDF 库")
    if not HAS_OPENPYXL:
        raise RuntimeError("PDF转Excel需要 openpyxl 库，请执行: pip install openpyxl")

    try:
        doc = fitz.open(pdf_path)
        wb = Workbook()
        default_ws = wb.active

        import re

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text")

            if not text.strip():
                continue

            if page_num == 0:
                ws = default_ws
                ws.title = "第1页"
            else:
                ws = wb.create_sheet(title=f"第{page_num+1}页")

            lines = text.split("\n")
            for row_idx, line in enumerate(lines, 1):
                cells = re.split(r'\t|\s{2,}', line.strip())
                for col_idx, cell_val in enumerate(cells, 1):
                    ws.cell(row=row_idx, column=col_idx, value=cell_val.strip())

            if progress_callback:
                progress_callback(page_num + 1, len(doc))

        doc.close()
        wb.save(output_path)
        return True

    except Exception as e:
        raise RuntimeError(f"PDF转Excel失败: {e}")
